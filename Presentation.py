# -*- coding: utf-8 -*-
"""
Created on Mon Jul 29 11:42:29 2019

@author: 
"""

import os, shutil, pptx, csv, datetime, matplotlib.pyplot as plt
from pptx.enum.text import PP_ALIGN
import numpy as np
import pandas as pd
from pptx.util import Inches, Pt, Cm

log_file = open('presentationcreationlog.txt', 'w')
log_file.write(str(datetime.datetime.now().strftime('%d %b %Y %A %H:%M:%S.%f')) + '\n')


#Размеры слайдов
slide_width=Cm(25.4)
slide_height=Cm(14.29)

prs = pptx.Presentation() #создание файла-презентации
prs.slide_height=slide_height
prs.slide_width=slide_width
log_file.write('Презентация создана\n\n')

#титульный слайд:
presentation_name='Снижение массы двигателя ТВ7-117СТ-01\v на основе\
 «цифрового двойника»'
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = presentation_name
title.text_frame.paragraphs[0].font.size = Pt(20)
title.text_frame.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(0x00, 0x00, 0xFF)
title.text_frame.paragraphs[0].font.bold = True
log_file.write('Титульный лист создан\n')

#  Поиск изображений (jpg, png) и копирование в ../pictures
log_file.write('Поиск изображений (jpg, png) и копирование в ../pictures \n')
imagelist=[]
values_dir = 'pictures'
if not os.path.exists(values_dir):
    os.makedirs(values_dir)
#Список файлов директории #directory    
#directory = 'pic'   #Список файлов директории ../pic
try:
    files = os.listdir()  #files = os.listdir(directory)
    imagelist = list(filter(lambda x: x.endswith('.jpg'), files))+\
    list(filter(lambda x: x.endswith('.png'), files))
    for image in imagelist:
        shutil.copy2(image, 'pictures/')  #shutil.copy2(directory+image, 'pictures/')
        log_file.write('Изображение "'+image+'" скопировано\n')
except: log_file.write('Изображений не найдено\n')



max_row=10  #Число строк в таблице 

def slide_with_table(max_row=max_row):
    """ 
    Создает слайд с таблицей. Для обращения к таблице глобальная\
    переменная table. Заполняет только заглавную строку, для заполнения \
    других строк функция table1(c)
    """
    global table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left=Inches(0.6)
    top=Inches(0.9)
    width=Cm(23)
    height=Cm(1)
    table = slide.shapes.add_table(rows=max_row+1, cols=4, left=left,\
                                   top=top, width=width, height=height)
    table.table.cell(0, 0).text='№'
    table.table.columns[0].width = Cm(1.5)
    table.table.cell(0, 1).text='Величина'
    table.table.columns[1].width = Cm(11)
    table.table.cell(0, 2).text='Значение'
    table.table.columns[2].width = Cm(5)
    table.table.cell(0, 3).text='Ед. изм.'
    table.table.columns[3].width = Cm(5)
    
    log_file.write('Создание слайда с таблицей\n')
    

       
    for i in range(4):
        table.table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(16)
    for i in [0,2,3]:   
        table.table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
def table1(c):
    
    log_file.write('Заполнение таблицы\n')
    
    slide_with_table(min(max_row,len(c)))
    row_idx=1
    n=0

    for c1 in c:
        table.table.cell(row_idx, 0).text=c1[0] 
        table.table.cell(row_idx, 1).text=c1[1] 
        table.table.cell(row_idx, 2).text=c1[2] 
        table.table.cell(row_idx, 3).text=dem_eng2rus(c1[5])
        n+=1
        for i in range(4):
            table.table.cell(row_idx, i).text_frame.paragraphs[0].font.size = Pt(16)
        for i in [0,2,3]:   
            table.table.cell(row_idx, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if row_idx>=max_row:
            slide_with_table(min(max_row,len(c)-n))
            row_idx=1
          
        else: row_idx+=1
    log_file.write('Таблица заполнена\n')
    #table.table.row(row_idx).height=1    
  
def dem_eng2rus(string):
    """
    Заменяет английскую размерность из списка eng на русскую из списка rus с\
    тем-же индексом
    """
    eng=['',   '[Pa]','[kg]','[s]', '[kg s^-1]', '[K]']
    rus=['[-]','[Па]','[кг]','[с]', '[кг/с]',    '[К]']
    try: dem_eng2rus=rus[eng.index(string)]
    except: dem_eng2rus=string
    return dem_eng2rus

def report_to_list(file='Report.txt'):
    """
    Читает файл file в двухмерный список для keypoint и построения таблицы
    """    
    #Открытие файла Report.txt
    f = open(file)
    line = f.readline()
    i=0
    c=[]
    while line:
        a=line[:len(line)-1].split('\t') #:len(line)-1  - Костыль чтобы убрать перенос на другую строку в конце строки
        b=[' ',' ',' ',' ',' ','',' ',' ']
        try:
            float(a[0])      
        except:
            pass
        else: 
  	        b[0]=str(float(a[0]))  #float
  	        b[1]=str(a[1])         #str
  	        try:b[2]=str(float(a[2]))  #float or str
  	        except:b[2]=str(a[2].split()[0])
  	        b[5]=str(a[3])[1:] #str  #  [1:]  - чтобы убрать пробел в начале ячейки с Ед. измерения
  	        c.append(b)
  	        i+=1
        line = f.readline() 
    f.close ()
    log_file.write('Файл "'+file+'" прочитан\n')
    return c
    

def report_to_list2(file='Report.txt', n=6):
    """
    Читает файл file в двухмерный список с числом столбцов n
    """
    #Открытие файла Report.txt
    f = open(file)
    line = f.readline()
    i=0
    c=[]
    while line:
        a=line[:len(line)-1].split('\t') #:len(line)-1  - Костыль чтобы убрать перенос на другую строку в конце строки
        b=[]
        try:
            float(a[0])      
        except:
            r=False
        else: r=True
            
        if r:
            for i in range(n):
                try: b.append(a[i])    
                except: b.append(' ')
            c.append(b)
            i+=1
        line = f.readline() 
    f.close ()
    log_file.write('Файл "'+file+'" прочитан\n')
    return c
    

def list_to_values(list_of_values=[' ',' ',' ',' ',' ','',' ',' ']):
    
    
    log_file.write('Запись keypoint\n')
    values_dir = 'values'
    if not os.path.exists(values_dir):
  	    os.makedirs(values_dir)
    values_file = values_dir+'/values.csv'
    f1=open(values_file, "w", newline='')
    writer = csv.writer(f1)
    writer.writerow(['#','Name','Value','Target','Condition',\
                     'Dimension','Description','Comments','']) 
    writer.writerows(list_of_values)   
    f1.close()
    log_file.write('Keypoint записаны в '+values_file+'\n')

     
     
def Slide_pictrue(picture_path0, height = Inches(2), printfilename=False, maxnpicslide=4):
    """Добавление изображения picture_path0 на слайд"""
    global npicslide
    global slide
    if npicslide==0:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = Inches(1+npicslide%(maxnpicslide/2)*(slide_width/Inches(1)-1)/(maxnpicslide/2))
    top = Inches(0.9+(npicslide>=maxnpicslide/2)*2.5)
    picture_path = 'pictures/'+picture_path0 #путь, где лежит изображение
    pic = slide.shapes.add_picture(picture_path, left, top,height=height)
    if printfilename:
        toptxt =  top-Inches(0.3)
        lefttxt=left
        widthtxt=Inches(4)
        txBox = slide.shapes.add_textbox(lefttxt, toptxt, widthtxt, Inches(0.5))
        tf = txBox.text_frame
        tf.text = picture_path0[0:len(picture_path0)-4]
        txBox.text_frame.paragraphs[0].font.size = Pt(12)
    npicslide+=1
    npicslide=npicslide%maxnpicslide   
    log_file.write('Изображение "'+picture_path0+'" добавлено на слайд\n')    

def plot_of_list(list):
    """
    Создает график из списка list. Столбцы списка: номер графика,\
     название графика и файла рисунка, x, y, [название оси x], [название оси y]
    """
    log_file.write('Создание графиков\n')
    fig, ax = plt.subplots()
    #fig.set_size_inches(10, 5)
    
    x=[]
    y=[]
    n=int(list[0][0])
    for list1 in list:
       
        if int(list1[0])!=n:

            ax.plot(x, y, 'o-')

            fig.savefig('pictures/'+filename, bbox_inches="tight")  
            #fig.show()
            Slide_pictrue(filename, height = Inches(2))
            #Slide_pictrue(filename)
            fig, ax = plt.subplots()
            
            x=[]
            y=[]            
            n=int(list1[0])
        #print('1  ',list1[2])
        x.append(float(list1[2]))
        y.append(float(list1[3])) 
        fig.set_size_inches(6, 3)
        ax.set_title(list1[1])
        filename=list1[1]+'.jpg'
        #print(str(list1[1]))
        try: ax.set_xlabel(r'{}'.format(list1[4]))
        except: pass
        try: ax.set_ylabel(r'{}'.format(list1[5]))
        except: pass

                
    ax.plot(x, y, 'o-')
    #plt.legend()
    fig.savefig('pictures/'+filename, bbox_inches="tight")
    #fig.show()
    Slide_pictrue(filename, height = Inches(2))
    log_file.write('Графики записаны\n')
        
    
list_of_values = report_to_list('Report.txt')
list_to_values(list_of_values)


#slide = prs.slides.add_slide(prs.slide_layouts[6])
npicslide=0
try: 
    for i in imagelist: Slide_pictrue(i, printfilename=True, maxnpicslide=6)
except: pass

#slide_with_table()  
table1(list_of_values)
#slide = prs.slides.add_slide(prs.slide_layouts[6])
npicslide=0
try:plot_of_list(report_to_list2(file='Report2.txt', n=6))
except: log_file.write('Таблиц для графиков не найдено\n')

def colontituls_and_page_numbers(presentation_name=presentation_name):
    log_file.write('Создание колонтитулов и номеров страниц\n')
    page_number=-1# Если =-1 - на первом (титульном) слайде номер слайда не ставится
    for slide in prs.slides:
        page_number+=1
        if page_number>0: 
            #Номер страницы
            txBox = slide.shapes.add_textbox(prs.slide_width-Cm(0.75), prs.slide_height-Cm(0.6), Cm(0.75), Cm(0.5))
            tf = txBox.text_frame
            tf.text = str(page_number)
            txBox.text_frame.paragraphs[0].font.size = Pt(10)
            #Нижний колонтитул
            txBox = slide.shapes.add_textbox(Cm(0), prs.slide_height-Cm(0.6), Cm(20), Cm(0.5))
            tf = txBox.text_frame
            tf.text = presentation_name
            txBox.text_frame.paragraphs[0].font.size = Pt(10)
            #Верхний колонтитул
            txBox = slide.shapes.add_textbox(prs.slide_width-Cm(3.3), Cm(0), Cm(3.3), Cm(0.9))
            tf = txBox.text_frame
            tf.text = 'CML-Bench\u2122'
            txBox.text_frame.paragraphs[0].font.size = Pt(14)
            txBox.text_frame.paragraphs[0].font.bold = 1
            txBox.text_frame.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(0x00, 0x00, 0xFF)
    log_file.write('Колонтитулы и номера слайдов до слайда №'+str(page_number)+' добавлены\n')
        
colontituls_and_page_numbers('Снижение массы двигателя ТВ7-117СТ-01 на основе\
 «цифрового двойника»')
        
prs.save('Presentation.pptx')
log_file.write('\nПрезентация записана в файл\n'+str(datetime.datetime.now().strftime('%H:%M:%S.%f')))

#Удаление файлов рисунков
"""try: 
    for i in imagelist: os.remove(i)
except: pass
"""

log_file.close ()
#Удаление папки     ../pic
#shutil.rmtree('pic/') # Временно закомментировано для отладки



